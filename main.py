import gradio as gr
import pyttsx3
from pptx import Presentation
import os
from io import BytesIO
from PIL import Image

# Initialize the text-to-speech engine
engine = pyttsx3.init()

# Create directories to store audio and images if they don't exist
audio_folder = "slide_audio_files"
image_folder = "slide_image_files"
os.makedirs(audio_folder, exist_ok=True)
os.makedirs(image_folder, exist_ok=True)

# Function to handle the action (Next only, no looping)
def handle_action(uploaded_file, current_slide):
    try:
        # Load the PowerPoint presentation
        presentation = Presentation(uploaded_file)

        # Get the current slide
        slide = presentation.slides[current_slide]

        # Extract text from the current slide
        slide_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text += shape.text + "\n"

        # Extract images from the slide (if any)
        image_path = None
        for shape in slide.shapes:
            # Check if the shape contains an image
            if hasattr(shape, "image") and shape.image:
                # Save the image to a file
                image_stream = BytesIO(shape.image.blob)
                image = Image.open(image_stream)
                image_path = os.path.join(image_folder, f"slide_image_{current_slide}.png")
                image.save(image_path)
                break  # Stop after saving the first image

        # Convert text to speech and save it as an audio file (only if text is present)
        audio_file = None
        if slide_text.strip():  # Only process if there is text
            audio_file = os.path.join(audio_folder, f"slide_audio_{current_slide}.mp3")
            engine.save_to_file(slide_text, audio_file)
            engine.runAndWait()

        # Move to the next slide
        current_slide += 1

        # Check if we've reached the last slide
        completion_message = ""
        if current_slide >= len(presentation.slides):
            current_slide = -1  # End of slides
            completion_message = "Complete!"

        # Ensure at least one output is returned (either text or image)
        if not slide_text.strip() and not image_path:
            slide_text = "No content found on this slide."

        return slide_text, audio_file, current_slide, image_path, completion_message

    except Exception as e:
        return f"Error: {str(e)}", None, current_slide, None, ""

# Function to define the interface
def interface():
    with gr.Blocks() as demo:
        # Upload PowerPoint file
        uploaded_file = gr.File(label="Upload PowerPoint File", type="filepath")

        # Current slide state
        current_slide = gr.State(0)

        # Output text for slide content
        slide_text = gr.Textbox(label="Slide Text", interactive=False)

        # Output audio for the slide
        audio_output = gr.Audio(label="Slide Audio", type="filepath")

        # Output image for the slide (if available)
        image_output = gr.Image(label="Slide Image", type="filepath")

        # Output completion message
        completion_message = gr.HTML()

        # Submit button to go to the next slide
        submit_btn = gr.Button("Next")
        submit_btn.click(
            handle_action,
            inputs=[uploaded_file, current_slide],
            outputs=[slide_text, audio_output, current_slide, image_output, completion_message],
        )

    return demo

# Launch the Gradio interface
if __name__ == "__main__":
    interface().launch()
